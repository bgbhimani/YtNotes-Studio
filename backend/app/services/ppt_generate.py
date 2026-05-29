# from pptx import Presentation
# from pptx.util import Pt


# # =========================================================
# # HELPER FUNCTION
# # =========================================================

# def get_layout_by_name(prs, layout_name):
#     """
#     Find layout by layout name
#     """
#     for layout in prs.slide_layouts:
#         if layout.name.lower() == layout_name.lower():
#             return layout

#     raise ValueError(f"Layout '{layout_name}' not found in template.")


# # =========================================================
# # MAIN PPT GENERATOR
# # =========================================================

# def generate_ppt(
#     slides_data,
#     output_file="generated_presentation.pptx",
#     template_path="app/src/Demo.pptx",
# ):
#     """
#     Generate PPT using your custom Slide Master layouts
#     """

#     # Load template
#     prs = Presentation(template_path)

#     # -----------------------------------------------------
#     # DEBUG: PRINT AVAILABLE LAYOUTS
#     # -----------------------------------------------------
#     print("\nAVAILABLE LAYOUTS:\n")

#     for i, layout in enumerate(prs.slide_layouts):
#         print(f"{i} -> {layout.name}")

#     print("\n----------------------------------\n")

#     # -----------------------------------------------------
#     # CREATE SLIDES
#     # -----------------------------------------------------
#     for slide_data in slides_data:

#         layout_type = slide_data.get("layout", "content").lower()

#         # -------------------------------------------------
#         # GET LAYOUT BY NAME
#         # -------------------------------------------------
#         slide_layout = get_layout_by_name(prs, layout_type)

#         # Create slide
#         slide = prs.slides.add_slide(slide_layout)

#         # -------------------------------------------------
#         # DEBUG PLACEHOLDERS
#         # -------------------------------------------------
#         print(f"\nSlide Layout: {layout_type}")

#         for placeholder in slide.placeholders:
#             print(
#                 f"IDX: {placeholder.placeholder_format.idx} "
#                 f"-> {placeholder.name}"
#             )

#         # -------------------------------------------------
#         # TITLE SLIDE
#         # -------------------------------------------------
#         if layout_type == "title":

#             # Title placeholder
#             if slide.shapes.title:
#                 slide.shapes.title.text = slide_data.get("title", "")

#             # Subtitle placeholder
#             if len(slide.placeholders) > 1:
#                 slide.placeholders[1].text = slide_data.get(
#                     "subtitle", ""
#                 )

#         # -------------------------------------------------
#         # CONTENT SLIDE
#         # -------------------------------------------------
#         elif layout_type == "content":

#             # Set title
#             if slide.shapes.title:
#                 slide.shapes.title.text = slide_data.get("title", "")

#             bullets = slide_data.get("bullets", [])

#             # Find content placeholder safely
#             content_placeholder = None

#             for placeholder in slide.placeholders:

#                 # BODY placeholder usually contains bullets
#                 if "content" in placeholder.name.lower() or \
#                    "text" in placeholder.name.lower():

#                     content_placeholder = placeholder
#                     break

#             # Fallback
#             if content_placeholder is None and len(slide.placeholders) > 1:
#                 content_placeholder = slide.placeholders[1]

#             # Add bullets
#             if content_placeholder:

#                 # text_frame = content_placeholder.text_frame
#                 # text_frame.clear()

#                 # for idx, bullet in enumerate(bullets):

#                 #     if idx == 0:
#                 #         p = text_frame.paragraphs[0]
#                 #     else:
#                 #         p = text_frame.add_paragraph()

#                 #     p.text = bullet
#                 #     p.font.size = Pt(20)
                
#                 text_frame = content_placeholder.text_frame

#                 # First bullet
#                 text_frame.paragraphs[0].text = bullets[0]

#                 # Remaining bullets
#                 for bullet in bullets[1:]:

#                     p = text_frame.add_paragraph()
#                     p.text = bullet

#         # -------------------------------------------------
#         # SECTION SLIDE
#         # -------------------------------------------------
#         elif layout_type == "section":

#             # Set title
#             if slide.shapes.title:
#                 slide.shapes.title.text = slide_data.get("title", "")

#             # Subtitle
#             if len(slide.placeholders) > 1:
#                 slide.placeholders[1].text = slide_data.get(
#                     "subtitle", ""
#                 )

#     # -----------------------------------------------------
#     # SAVE
#     # -----------------------------------------------------
#     prs.save(output_file)

#     print(f"\nPPT saved successfully: {output_file}")



# # # =========================================================
# # # EXAMPLE DATA
# # # =========================================================

# # slides_data = [
# #     {
# #         "layout": "title",
# #         "title": "Future of AI",
# #         "subtitle": "Generated from YouTube Transcript"
# #     },
# #     {
# #         "layout": "section",
# #         "title": "Introduction",
# #         "subtitle": "Understanding AI"
# #     },
# #     {
# #         "layout": "content",
# #         "title": "What is AI?",
# #         "bullets": [
# #             "Artificial Intelligence simulates human thinking",
# #             "Used in automation",
# #             "Improves productivity",
# #             "Supports decision making"
# #         ]
# #     }
# # ]
# # generate_ppt(slides_data)


from pptx import Presentation
from pptx.util import Pt


# =========================================================
# HELPER FUNCTION
# =========================================================

def get_layout_by_name(prs, layout_name):
    """
    Find slide layout by layout name
    """

    for layout in prs.slide_layouts:

        if layout.name.lower() == layout_name.lower():
            return layout

    raise ValueError(
        f"Layout '{layout_name}' not found in template."
    )


# =========================================================
# MAIN PPT GENERATOR
# =========================================================

def generate_ppt(
    slides_data,
    output_file="generated_presentation.pptx",
    template_path="app/src/Demo.pptx",
):
    """
    Generate PowerPoint presentation using
    custom Slide Master layouts.
    Make Sure to Make around 5-15 Slides Only.
    Try to keep the content concise for best results.

    Supported layouts:
    - title
    - section
    - content
    """

    # =====================================================
    # LOAD TEMPLATE
    # =====================================================

    prs = Presentation(template_path)

    # =====================================================
    # CREATE SLIDES
    # =====================================================

    for slide_data in slides_data:

        layout_type = slide_data.get(
            "layout",
            "content"
        ).lower()

        # -------------------------------------------------
        # GET LAYOUT
        # -------------------------------------------------

        slide_layout = get_layout_by_name(
            prs,
            layout_type
        )

        # -------------------------------------------------
        # CREATE SLIDE
        # -------------------------------------------------

        slide = prs.slides.add_slide(slide_layout)

        # =================================================
        # TITLE SLIDE
        # =================================================

        if layout_type == "title":

            # 🔹 Title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get(
                    "title",
                    ""
                )

            # 🔹 Subtitle
            if len(slide.placeholders) > 1:

                slide.placeholders[1].text = slide_data.get(
                    "subtitle",
                    ""
                )

        # =================================================
        # SECTION SLIDE
        # =================================================

        elif layout_type == "section":

            # 🔹 Title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get(
                    "title",
                    ""
                )

            # 🔹 Subtitle
            if len(slide.placeholders) > 1:

                slide.placeholders[1].text = slide_data.get(
                    "subtitle",
                    ""
                )

        # =================================================
        # CONTENT SLIDE
        # =================================================

        elif layout_type == "content":

            # 🔹 Title
            if slide.shapes.title:

                slide.shapes.title.text = slide_data.get(
                    "title",
                    ""
                )

            bullets = slide_data.get("bullets", [])

            # -------------------------------------------------
            # FIND CONTENT PLACEHOLDER
            # -------------------------------------------------

            content_placeholder = None

            for placeholder in slide.placeholders:

                name = placeholder.name.lower()

                if (
                    "content" in name
                    or "text" in name
                    or "body" in name
                ):

                    content_placeholder = placeholder
                    break

            # 🔹 fallback
            if (
                content_placeholder is None
                and len(slide.placeholders) > 1
            ):
                content_placeholder = slide.placeholders[1]

            # -------------------------------------------------
            # ADD BULLETS
            # -------------------------------------------------

            if content_placeholder and bullets:

                text_frame = content_placeholder.text_frame

                # Clear existing placeholder text
                text_frame.clear()

                # 🔹 First bullet
                first_para = text_frame.paragraphs[0]

                first_para.text = bullets[0]
                first_para.font.size = Pt(20)

                # 🔹 Remaining bullets
                for bullet in bullets[1:]:

                    p = text_frame.add_paragraph()

                    p.text = bullet
                    p.font.size = Pt(20)

    # =====================================================
    # SAVE PPT
    # =====================================================

    prs.save(output_file)

    print(f"\n✅ PPT saved successfully: {output_file}")